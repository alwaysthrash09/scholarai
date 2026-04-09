import streamlit as st
import anthropic
import tempfile
import os
import io
from pathlib import Path
from datetime import datetime

# ─────────────────────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="ScholarAI — BC MBA",
    page_icon="🦅",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ─────────────────────────────────────────────────────────────────────────────
COURSES = ["Data Analytics", "Finance", "Strategy",
           "Marketing", "Operations", "Economics"]

COURSE_COLORS = {
    "Data Analytics": "#4FC9A4",
    "Finance":        "#F7C46A",
    "Strategy":       "#7C6AF7",
    "Marketing":      "#F76A6A",
    "Operations":     "#6AB8F7",
    "Economics":      "#C46AF7",
}

COURSE_ICONS = {
    "Data Analytics": "📊",
    "Finance":        "💰",
    "Strategy":       "♟️",
    "Marketing":      "📣",
    "Operations":     "⚙️",
    "Economics":      "📈",
}

# ─────────────────────────────────────────────────────────────────────────────
# Custom CSS
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=DM+Sans:wght@300;400;500;700&family=DM+Serif+Display&display=swap');

    html, body, [class*="css"] { font-family: 'DM Sans', sans-serif; }

    .main { background-color: #0A0A10; }
    .block-container { padding-top: 1.5rem; padding-bottom: 2rem; }

    .app-header {
        background: #12121A;
        border-radius: 16px;
        padding: 22px 28px;
        margin-bottom: 20px;
        border: 1px solid #2A2A3A;
        display: flex;
        align-items: center;
        justify-content: space-between;
    }
    .app-title {
        font-family: 'DM Serif Display', serif;
        font-size: 2rem;
        color: #E8E8F0;
        margin: 0;
    }
    .app-subtitle { color: #5A5A72; font-size: 0.9rem; margin-top: 2px; }

    .course-pill {
        display: inline-block;
        padding: 6px 16px;
        border-radius: 20px;
        font-size: 0.85rem;
        font-weight: 700;
        margin-right: 6px;
        margin-bottom: 6px;
        cursor: pointer;
        border: 2px solid transparent;
    }

    .section-label {
        color: #5A5A72;
        font-size: 0.7rem;
        font-weight: 700;
        text-transform: uppercase;
        letter-spacing: 1.5px;
        margin-bottom: 8px;
        margin-top: 16px;
    }

    .material-item {
        background: #16161E;
        border-radius: 8px;
        padding: 8px 12px;
        margin-bottom: 6px;
        border: 1px solid #2A2A3A;
        font-size: 0.85rem;
        color: #C0C0D8;
    }

    .chat-container {
        background: #0A0A10;
        border-radius: 12px;
        border: 1px solid #2A2A3A;
        min-height: 400px;
        padding: 16px;
        margin-bottom: 12px;
    }

    .quick-action-grid {
        display: grid;
        grid-template-columns: 1fr 1fr;
        gap: 8px;
        margin-bottom: 16px;
    }

    /* Sidebar */
    [data-testid="stSidebar"] { background-color: #12121A !important; }
    [data-testid="stSidebar"] h1,
    [data-testid="stSidebar"] h2,
    [data-testid="stSidebar"] h3,
    [data-testid="stSidebar"] p,
    [data-testid="stSidebar"] label { color: #E8E8F0 !important; }

    /* Buttons */
    .stButton button {
        border-radius: 8px;
        font-weight: 600;
        font-family: 'DM Sans', sans-serif;
        border: none;
    }
    .stButton button[kind="primary"] {
        background: #7C6AF7 !important;
        color: white !important;
    }
    .stDownloadButton button {
        background: #4FC9A4 !important;
        color: #0A0A10 !important;
        border: none !important;
        border-radius: 8px !important;
        font-weight: 700 !important;
    }

    /* File uploader */
    [data-testid="stFileUploader"] {
        background: #16161E;
        border-radius: 12px;
        border: 2px dashed #2A2A3A;
        padding: 8px;
    }

    /* Text input */
    .stTextInput input, .stTextArea textarea {
        background: #16161E !important;
        color: #E8E8F0 !important;
        border: 1px solid #2A2A3A !important;
        border-radius: 8px !important;
    }

    /* Select */
    .stSelectbox > div > div {
        background: #16161E !important;
        color: #E8E8F0 !important;
        border: 1px solid #2A2A3A !important;
    }

    /* Tabs */
    .stTabs [data-baseweb="tab-list"] {
        background: #12121A;
        border-radius: 10px;
        padding: 4px;
        gap: 4px;
    }
    .stTabs [data-baseweb="tab"] {
        background: transparent;
        color: #5A5A72;
        border-radius: 8px;
        font-weight: 600;
        padding: 8px 16px;
    }
    .stTabs [aria-selected="true"] {
        background: #1E1E2A !important;
        color: #E8E8F0 !important;
    }

    hr { border-color: #2A2A3A; }
    .stAlert { border-radius: 10px; }
    .stSuccess { background: #0A2A1A !important; color: #4FC9A4 !important; }
    .stInfo    { background: #0A0A2A !important; color: #7C6AF7 !important; }
</style>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────────────────────────────────────────
# Session state init
def init_state():
    if "current_course" not in st.session_state:
        st.session_state.current_course = COURSES[0]
    if "materials" not in st.session_state:
        st.session_state.materials = {c: [] for c in COURSES}
    if "histories" not in st.session_state:
        st.session_state.histories = {c: [] for c in COURSES}
    if "api_key" not in st.session_state:
        st.session_state.api_key = ""
    if "student_name" not in st.session_state:
        st.session_state.student_name = ""

init_state()

# ─────────────────────────────────────────────────────────────────────────────
def extract_pdf(file_bytes):
    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(file_bytes))
        return "\n".join(p.extract_text() or "" for p in reader.pages)[:6000]
    except Exception as e:
        return f"[Could not read PDF: {e}]"

def extract_pptx(file_bytes):
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)[:6000]
    except Exception as e:
        return f"[Could not read PPTX: {e}]"

def get_context(course):
    mats = st.session_state.materials.get(course, [])
    parts = []
    for m in mats:
        parts.append(f"=== {m['name']} ===\n{m['content'][:3000]}")
    return "\n\n".join(parts)

def chat_with_claude(user_message, course, system_override=None):
    key = st.session_state.api_key
    if not key:
        return "⚠️ Please enter your Anthropic API key in the sidebar."

    context = get_context(course)
    name = st.session_state.student_name
    student_note = f" The student's name is {name}." if name else ""

    system = system_override or f"""You are ScholarAI, a dedicated MBA study tutor for a Boston College MBA part-time student who also works full-time.{student_note}

Current course: {course}

Your personality:
- Conversational and encouraging like a great personal tutor
- Thorough but respects the student's limited time
- Uses plain English — avoids unnecessary jargon
- Connects concepts to real-world business applications
- Aware this is an MBA student — assumes professional maturity

{"COURSE MATERIALS:" + chr(10) + context if context else "No materials uploaded yet for this course. Encourage the student to upload their materials."}

Structure longer responses with clear headers and bullet points for easy scanning.
End responses with a follow-up question or suggestion to keep learning going."""

    history = st.session_state.histories[course][-10:]
    messages = list(history) + [{"role": "user", "content": user_message}]

    try:
        client = anthropic.Anthropic(api_key=key)
        response = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=2000,
            system=system,
            messages=messages)
        return response.content[0].text
    except Exception as e:
        return f"❌ Error: {e}"

# ─────────────────────────────────────────────────────────────────────────────
# SIDEBAR
with st.sidebar:
    st.markdown("## 🦅 ScholarAI")
    st.markdown("*BC MBA Study Agent*")
    st.markdown("---")

    st.markdown("### ⚙️ Settings")
    api_key_input = st.text_input(
        "Anthropic API Key",
        type="password",
        placeholder="sk-ant-...",
        value=st.session_state.api_key,
        help="Get your key at console.anthropic.com")
    if api_key_input:
        st.session_state.api_key = api_key_input

    name_input = st.text_input(
        "Your Name",
        placeholder="e.g. Howard",
        value=st.session_state.student_name)
    if name_input:
        st.session_state.student_name = name_input

    st.markdown("---")
    st.markdown("### 📚 Select Course")
    for course in COURSES:
        color = COURSE_COLORS[course]
        icon  = COURSE_ICONS[course]
        is_active = st.session_state.current_course == course
        if st.button(
            f"{icon}  {course}",
            key=f"course_btn_{course}",
            type="primary" if is_active else "secondary",
            use_container_width=True):
            st.session_state.current_course = course
            st.rerun()

    st.markdown("---")
    st.markdown("### 📎 Add Materials")
    st.markdown(f"*Adding to: **{st.session_state.current_course}***")

    uploaded_file = st.file_uploader(
        "Upload File",
        type=["pdf", "pptx", "txt", "md"],
        help="PDF, PowerPoint, or text files",
        key="file_uploader")

    if uploaded_file:
        with st.spinner("Reading file..."):
            suffix = Path(uploaded_file.name).suffix.lower()
            file_bytes = uploaded_file.read()
            if suffix == ".pdf":
                content = extract_pdf(file_bytes)
            elif suffix == ".pptx":
                content = extract_pptx(file_bytes)
            else:
                content = file_bytes.decode("utf-8", errors="ignore")[:6000]

            course = st.session_state.current_course
            # Avoid duplicates
            existing_names = [m["name"] for m in st.session_state.materials[course]]
            if uploaded_file.name not in existing_names:
                st.session_state.materials[course].append({
                    "name": uploaded_file.name,
                    "type": suffix,
                    "content": content,
                    "added": datetime.now().strftime("%b %d %I:%M %p")
                })
                st.success(f"✅ Added: {uploaded_file.name}")
            else:
                st.info("Already added!")

    url_input = st.text_input("Or paste a URL", placeholder="https://...")
    if st.button("+ Add URL", use_container_width=True) and url_input:
        with st.spinner("Fetching URL..."):
            try:
                import urllib.request
                req = urllib.request.Request(
                    url_input, headers={"User-Agent": "Mozilla/5.0"})
                with urllib.request.urlopen(req, timeout=10) as r:
                    html = r.read().decode("utf-8", errors="ignore")
                import re
                html = re.sub(r"<[^>]+>", " ", html)
                html = re.sub(r"\s+", " ", html).strip()[:4000]
                course = st.session_state.current_course
                st.session_state.materials[course].append({
                    "name": url_input[:50],
                    "type": "url",
                    "content": html,
                    "added": datetime.now().strftime("%b %d %I:%M %p")
                })
                st.success("✅ URL added!")
            except Exception as e:
                st.error(f"Could not fetch: {e}")

    # StudyScribe transcript upload
    st.markdown("---")
    transcript_file = st.file_uploader(
        "🎙 Import StudyScribe Transcript",
        type=["txt", "md"],
        key="transcript_uploader",
        help="Upload a .txt or .md file exported from StudyScribe")
    if transcript_file:
        content = transcript_file.read().decode("utf-8", errors="ignore")
        course  = st.session_state.current_course
        existing_names = [m["name"] for m in st.session_state.materials[course]]
        if transcript_file.name not in existing_names:
            st.session_state.materials[course].append({
                "name": f"[Transcript] {transcript_file.name}",
                "type": "transcript",
                "content": content[:6000],
                "added": datetime.now().strftime("%b %d %I:%M %p")
            })
            st.success(f"✅ Transcript added!")

    st.markdown("---")
    st.markdown(
        "<div style='color:#3A3A5A; font-size:0.78rem;'>"
        "ScholarAI · BC MBA<br>Powered by Claude AI"
        "</div>", unsafe_allow_html=True)

# ─────────────────────────────────────────────────────────────────────────────
# MAIN CONTENT
course = st.session_state.current_course
color  = COURSE_COLORS[course]
icon   = COURSE_ICONS[course]

# Header
st.markdown(f"""
<div class="app-header">
    <div>
        <div class="app-title">🦅 ScholarAI</div>
        <div class="app-subtitle">Boston College MBA · Part-time Student Study Agent</div>
    </div>
    <div style="text-align:right;">
        <div style="font-size:1.8rem;">{icon}</div>
        <div style="color:{color}; font-weight:700; font-size:1rem;">{course}</div>
    </div>
</div>
""", unsafe_allow_html=True)

# ── Two column layout ────────────────────────────────────────────────────────
left_col, right_col = st.columns([1, 2], gap="large")

# ── LEFT: Materials + Quick Actions ─────────────────────────────────────────
with left_col:
    # Materials list
    st.markdown(f'<div class="section-label">📎 {course} Materials</div>',
                unsafe_allow_html=True)

    mats = st.session_state.materials.get(course, [])
    if mats:
        for i, m in enumerate(mats):
            type_icon = {"pdf":"📄","pptx":"📊",".txt":"📝",
                         "transcript":"🎙","url":"🌐"}.get(m["type"], "📎")
            col_a, col_b = st.columns([5, 1])
            with col_a:
                st.markdown(
                    f'<div class="material-item">{type_icon} {m["name"][:35]}<br>'
                    f'<span style="color:#3A3A5A;font-size:0.75rem;">{m["added"]}</span></div>',
                    unsafe_allow_html=True)
            with col_b:
                if st.button("✕", key=f"remove_{course}_{i}", help="Remove"):
                    st.session_state.materials[course].pop(i)
                    st.rerun()
    else:
        st.markdown(
            '<div style="color:#3A3A5A; font-size:0.85rem; padding:12px;">No materials yet — upload files in the sidebar!</div>',
            unsafe_allow_html=True)

    st.markdown("---")

    # Quick Actions
    st.markdown('<div class="section-label">⚡ Quick Actions</div>',
                unsafe_allow_html=True)

    actions = [
        ("📋 Module Report",      f"Generate a detailed weekly module report for my {course} course based on the materials I've uploaded. Include key concepts, frameworks, real-world applications, and what to focus on for exams."),
        ("📚 Study Guide",        f"Create a comprehensive study guide for {course} based on my uploaded materials. Organize by topic with key terms, core concepts, and important things to remember."),
        ("❓ Quiz Me",            f"Quiz me on {course}! Give me 5 MBA-level questions based on my course materials — mix conceptual and applied. After I answer, give me detailed feedback."),
        ("📅 Study Plan",         f"Based on my {course} materials, help me create a realistic study plan. I'm a part-time MBA student working full-time so time is limited. Prioritize important topics and suggest how many hours to allocate."),
        ("✏️ Assignment Help",    f"I need help with a {course} assignment. Ask me about the assignment details so you can help me structure my approach and address all requirements."),
        ("💡 Explain a Concept",  f"Ask me which concept from {course} I'd like explained in simple, plain English with a real-world business example."),
    ]

    for label, prompt in actions:
        if st.button(label, key=f"action_{label}_{course}", use_container_width=True):
            if not st.session_state.api_key:
                st.warning("Please enter your API key in the sidebar first.")
            else:
                with st.spinner("ScholarAI is thinking..."):
                    reply = chat_with_claude(prompt, course)
                    st.session_state.histories[course].append(
                        {"role": "user", "content": prompt})
                    st.session_state.histories[course].append(
                        {"role": "assistant", "content": reply})
                    st.rerun()

    st.markdown("---")
    if st.button("🗑️ Clear Chat History", use_container_width=True):
        st.session_state.histories[course] = []
        st.rerun()

# ── RIGHT: Chat ──────────────────────────────────────────────────────────────
with right_col:
    st.markdown(f'<div class="section-label">💬 Chat — {course}</div>',
                unsafe_allow_html=True)

    # Welcome message if empty
    history = st.session_state.histories[course]
    if not history:
        name = st.session_state.student_name
        greeting = f"Hi {name}! " if name else "Hi! "
        st.info(
            f"{greeting}I'm your ScholarAI tutor for **{course}**. "
            f"Upload your materials in the sidebar, then ask me anything or use the Quick Actions on the left. "
            f"I'm here to help you ace your BC MBA! 🦅")

    # Chat history
    for msg in history:
        role = msg["role"]
        with st.chat_message(role, avatar="🦅" if role == "assistant" else "👤"):
            st.markdown(msg["content"])

    # Download notes button if there's content
    if history:
        all_content = "\n\n---\n\n".join(
            [f"**{'You' if m['role']=='user' else 'ScholarAI'}:** {m['content']}"
             for m in history])
        ts = datetime.now().strftime("%Y%m%d_%H%M")
        st.download_button(
            f"💾 Download {course} Notes",
            data=all_content,
            file_name=f"ScholarAI_{course.replace(' ','_')}_{ts}.md",
            mime="text/markdown",
            use_container_width=True)

    # Chat input
    user_input = st.chat_input(
        f"Ask anything about {course}...")
    if user_input:
        if not st.session_state.api_key:
            st.warning("Please enter your Anthropic API key in the sidebar.")
        else:
            with st.chat_message("user", avatar="👤"):
                st.markdown(user_input)

            with st.chat_message("assistant", avatar="🦅"):
                with st.spinner("Thinking..."):
                    reply = chat_with_claude(user_input, course)
                    st.markdown(reply)

            st.session_state.histories[course].append(
                {"role": "user", "content": user_input})
            st.session_state.histories[course].append(
                {"role": "assistant", "content": reply})
            st.rerun()
